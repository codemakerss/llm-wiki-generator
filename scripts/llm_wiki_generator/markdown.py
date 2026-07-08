from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader


SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".txt", ".md", ".markdown"}


@dataclass
class MarkdownDocument:
    source_path: Path
    title: str
    markdown: str
    extension: str


def assert_supported(source_path: Path) -> None:
    if source_path.suffix.lower() not in SUPPORTED_EXTENSIONS:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise ValueError(f"Unsupported file type: {source_path.suffix}. Supported: {supported}")


def convert_docx(source_path: Path) -> tuple[str, str]:
    document = DocxDocument(source_path)
    lines = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    title = lines[0] if lines else source_path.stem
    return title, "\n\n".join(lines)


def convert_pptx(source_path: Path) -> tuple[str, str]:
    presentation = Presentation(source_path)
    slides: list[str] = []
    title = source_path.stem
    for index, slide in enumerate(presentation.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                texts.append(text)
        if texts:
            if index == 1:
                title = texts[0].splitlines()[0].strip() or title
            slides.append(f"## Slide {index}\n\n" + "\n\n".join(texts))
    return title, "\n\n".join(slides)


def convert_xlsx(source_path: Path) -> tuple[str, str]:
    workbook = load_workbook(source_path, data_only=True)
    sections: list[str] = []
    for sheet in workbook.worksheets:
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            values = [str(value).strip() for value in row if value is not None and str(value).strip()]
            if values:
                rows.append(" | ".join(values))
        if rows:
            sections.append(f"## {sheet.title}\n\n" + "\n".join(rows))
    return source_path.stem, "\n\n".join(sections)


def convert_pdf(source_path: Path) -> tuple[str, str]:
    reader = PdfReader(str(source_path))
    pages: list[str] = []
    for index, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if text:
            pages.append(f"## Page {index}\n\n{text}")
    return source_path.stem, "\n\n".join(pages)


def convert_to_markdown(source_path: Path) -> MarkdownDocument:
    source_path = source_path.resolve()
    assert_supported(source_path)
    if not source_path.exists():
        raise FileNotFoundError(source_path)

    suffix = source_path.suffix.lower()
    if suffix in {".txt", ".md", ".markdown"}:
        markdown = source_path.read_text(encoding="utf-8")
        title = source_path.stem
    elif suffix == ".docx":
        title, markdown = convert_docx(source_path)
    elif suffix == ".pptx":
        title, markdown = convert_pptx(source_path)
    elif suffix == ".xlsx":
        title, markdown = convert_xlsx(source_path)
    else:
        title, markdown = convert_pdf(source_path)

    return MarkdownDocument(
        source_path=source_path,
        title=title or source_path.stem,
        markdown=markdown.strip() + "\n",
        extension=source_path.suffix.lower(),
    )
