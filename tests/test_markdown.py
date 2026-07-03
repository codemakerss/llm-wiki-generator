from pathlib import Path

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from reportlab.pdfgen import canvas

from llm_wiki_generator.markdown import SUPPORTED_EXTENSIONS, assert_supported, convert_to_markdown


def test_supported_extensions_are_fixed() -> None:
    assert SUPPORTED_EXTENSIONS == {".pdf", ".docx", ".pptx", ".xlsx", ".txt"}


def test_convert_txt_to_markdown(tmp_path: Path) -> None:
    source = tmp_path / "sample.txt"
    source.write_text("hello\nworld\n", encoding="utf-8")

    document = convert_to_markdown(source)

    assert document.title == "sample"
    assert "hello" in document.markdown
    assert "world" in document.markdown


def test_unsupported_extension_raises(tmp_path: Path) -> None:
    source = tmp_path / "sample.md"
    source.write_text("# nope", encoding="utf-8")

    with pytest.raises(ValueError):
        assert_supported(source)


def test_converts_all_supported_formats(tmp_path: Path) -> None:
    files: list[Path] = []

    txt_path = tmp_path / "sample.txt"
    txt_path.write_text("business background\nmetrics\n", encoding="utf-8")
    files.append(txt_path)

    docx_path = tmp_path / "sample.docx"
    doc = Document()
    doc.add_heading("Team History", 0)
    doc.add_paragraph("Historical PRD notes.")
    doc.save(docx_path)
    files.append(docx_path)

    pptx_path = tmp_path / "sample.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Industry Practice"
    slide.placeholders[1].text = "PRD structure\nReview flow"
    prs.save(pptx_path)
    files.append(pptx_path)

    xlsx_path = tmp_path / "sample.xlsx"
    wb = Workbook()
    ws = wb.active
    ws["A1"] = "Metric"
    ws["B1"] = "85%"
    wb.save(xlsx_path)
    files.append(xlsx_path)

    pdf_path = tmp_path / "sample.pdf"
    pdf = canvas.Canvas(str(pdf_path))
    pdf.drawString(100, 750, "Risk control")
    pdf.drawString(100, 730, "P99 < 50ms")
    pdf.save()
    files.append(pdf_path)

    converted = [convert_to_markdown(path) for path in files]

    assert len(converted) == 5
    assert all(document.markdown.strip() for document in converted)
